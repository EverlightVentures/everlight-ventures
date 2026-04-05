"""
PowerPoint Generator - Branded Everlight Ventures presentations.
Uses python-pptx for professional deck generation.
"""
import os
from datetime import datetime
from pathlib import Path

OUTPUT_DIR = Path(os.environ.get("DELIVERABLES_DIR", "/tmp/hive_deliverables"))
OUTPUT_DIR.mkdir(parents=True, exist_ok=True)


def generate_pptx(title: str, slides: list[dict], subtitle: str = "",
                  author: str = "Everlight Ventures") -> str:
    """
    Generate a branded PowerPoint presentation.

    Args:
        title: Presentation title (first slide)
        slides: List of {"title": str, "body": str, "bullets": list[str] (optional)}
        subtitle: Subtitle for title slide
        author: Author metadata

    Returns:
        Path to generated .pptx file
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        return _fallback_text(title, slides, subtitle)

    prs = Presentation()
    prs.core_properties.author = author
    prs.core_properties.title = title

    # Slide dimensions (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    GOLD = RGBColor(0xD4, 0xA0, 0x17)
    DARK = RGBColor(0x1A, 0x1A, 0x1A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY = RGBColor(0x88, 0x88, 0x88)

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Dark background
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = DARK

    # Gold accent bar at top
    from pptx.util import Inches as In
    shape = slide.shapes.add_shape(
        1, 0, 0, prs.slide_width, Inches(0.15)  # Rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = GOLD
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(1.5))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    txBox2 = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11), Inches(0.8))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = subtitle or f"{author} | {datetime.now().strftime('%B %Y')}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = GOLD
    p2.alignment = PP_ALIGN.LEFT

    # --- Content Slides ---
    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # White background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE

        # Gold bar
        shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(0.08))
        shape.fill.solid()
        shape.fill.fore_color.rgb = GOLD
        shape.line.fill.background()

        # Slide title
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = DARK

        # Body text or bullets
        content_top = Inches(1.5)
        if slide_data.get("bullets"):
            txBox = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, bullet in enumerate(slide_data["bullets"]):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = DARK
                p.space_after = Pt(12)
                p.level = 0

        elif slide_data.get("body"):
            txBox = slide.shapes.add_textbox(Inches(0.8), content_top, Inches(11), Inches(5))
            tf = txBox.text_frame
            tf.word_wrap = True
            for i, line in enumerate(slide_data["body"].split("\n")):
                if not line.strip():
                    continue
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = line
                p.font.size = Pt(16)
                p.font.color.rgb = DARK
                p.space_after = Pt(8)

        # Footer
        txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(11), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"Everlight Ventures | everlightventures.io | Confidential"
        p.font.size = Pt(9)
        p.font.color.rgb = GRAY

    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M')}.pptx"
    filepath = OUTPUT_DIR / filename
    prs.save(str(filepath))
    return str(filepath)


def _fallback_text(title, slides, subtitle):
    """Fallback when python-pptx is not installed."""
    filename = f"{title.replace(' ', '_')}_{datetime.now().strftime('%Y%m%d_%H%M')}.txt"
    filepath = OUTPUT_DIR / filename
    lines = [f"PRESENTATION: {title}", subtitle or "", "=" * 40, ""]
    for s in slides:
        lines.append(f"\n--- SLIDE: {s.get('title', '')} ---")
        if s.get("body"):
            lines.append(s["body"])
        if s.get("bullets"):
            for b in s["bullets"]:
                lines.append(f"  * {b}")
    filepath.write_text("\n".join(lines))
    return str(filepath)
